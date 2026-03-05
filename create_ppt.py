from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 宽屏
prs.slide_height = Inches(7.5)

# 定义颜色主题 - 专业蓝色科技风
PRIMARY_BLUE = RgbColor(0, 82, 147)      # 深蓝
ACCENT_BLUE = RgbColor(0, 150, 199)      # 亮蓝
DARK_NAVY = RgbColor(10, 35, 66)         # 藏青
LIGHT_BLUE = RgbColor(230, 245, 255)     # 浅蓝背景
WHITE = RgbColor(255, 255, 255)
DARK_GRAY = RgbColor(51, 51, 51)

def add_gradient_background(slide, color1, color2):
    """添加渐变背景"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color1
    background.line.fill.background()
    # 将背景移到最底层
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_title_slide(prs, title, subtitle):
    """创建标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加深色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_NAVY
    bg.line.fill.background()
    
    # 添加装饰线条
    line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.2), Inches(2), Pt(4))
    line1.fill.solid()
    line1.fill.fore_color.rgb = ACCENT_BLUE
    line1.line.fill.background()
    
    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.333), Inches(3.2), Inches(2), Pt(4))
    line2.fill.solid()
    line2.fill.fore_color.rgb = ACCENT_BLUE
    line2.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # 底部装饰
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.8), prs.slide_width, Inches(0.7))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = PRIMARY_BLUE
    bottom_bar.line.fill.background()
    
    return slide

def add_content_slide(prs, title, content_items, has_icon=False):
    """创建内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 白色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # 顶部蓝色条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PRIMARY_BLUE
    top_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 内容
    y_pos = 1.6
    for item in content_items:
        if isinstance(item, tuple):  # (标题, 内容)
            item_title, item_content = item
            # 小标题
            box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(11.5), Inches(0.5))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = "▶ " + item_title
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_BLUE
            
            # 内容
            y_pos += 0.5
            box = slide.shapes.add_textbox(Inches(1.3), Inches(y_pos), Inches(11), Inches(0.8))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item_content
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.line_spacing = 1.3
            
            y_pos += 1.0
        else:
            box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(11.5), Inches(0.5))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = "• " + item
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
            y_pos += 0.6
    
    return slide

def add_product_slide(prs, title, products):
    """创建产品展示页 - 卡片式布局"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 浅色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BLUE
    bg.line.fill.background()
    
    # 标题区
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = DARK_NAVY
    title_bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 产品卡片布局 (2x2)
    card_width = Inches(5.8)
    card_height = Inches(2.6)
    positions = [(Inches(0.7), Inches(1.7)), (Inches(6.9), Inches(1.7)),
                 (Inches(0.7), Inches(4.5)), (Inches(6.9), Inches(4.5))]
    
    for i, (prod_name, prod_desc) in enumerate(products):
        if i >= 4:
            break
        x, y = positions[i]
        
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = ACCENT_BLUE
        card.line.width = Pt(2)
        card.shadow.inherit = False
        
        # 卡片顶部色条
        card_top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_width, Inches(0.5))
        card_top.fill.solid()
        card_top.fill.fore_color.rgb = PRIMARY_BLUE
        card_top.line.fill.background()
        
        # 产品名称
        name_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.08), card_width - Inches(0.4), Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = prod_name
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # 产品描述
        desc_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.7), card_width - Inches(0.6), Inches(1.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = prod_desc
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.line_spacing = 1.2
    
    return slide

def add_advantage_slide(prs):
    """创建核心优势页 - 圆形图标布局"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 白色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "核心竞争优势"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_NAVY
    p.alignment = PP_ALIGN.CENTER
    
    # 标题下划线
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(1.3), Inches(2.333), Pt(4))
    underline.fill.solid()
    underline.fill.fore_color.rgb = ACCENT_BLUE
    underline.line.fill.background()
    
    # 8个优势点 (2行4列)
    advantages = [
        ("15年", "行业沉淀经验"),
        ("80%", "技术+服务人员占比"),
        ("1200+", "成功客户案例"),
        ("7×24", "小时客服在线"),
        ("首创", "多组织+多仓+联盟模式"),
        ("百万票次", "年订单系统考验"),
        ("定制开发", "灵活可拓展"),
        ("1亿+", "商品池资源"),
    ]
    
    positions = [
        (Inches(0.8), Inches(1.8)), (Inches(3.6), Inches(1.8)), (Inches(6.4), Inches(1.8)), (Inches(9.2), Inches(1.8)),
        (Inches(0.8), Inches(4.5)), (Inches(3.6), Inches(4.5)), (Inches(6.4), Inches(4.5)), (Inches(9.2), Inches(4.5)),
    ]
    
    for i, (num, desc) in enumerate(advantages):
        x, y = positions[i]
        
        # 圆形背景
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y, Inches(1.5), Inches(1.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY_BLUE
        circle.line.fill.background()
        
        # 数字
        num_box = slide.shapes.add_textbox(x, y + Inches(0.35), Inches(2.1), Inches(0.8))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # 描述
        desc_box = slide.shapes.add_textbox(x, y + Inches(1.6), Inches(2.1), Inches(0.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_client_slide(prs):
    """创建客户案例页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 深色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_NAVY
    bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "代表客户"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 客户名称展示
    clients = [
        "中欧联运", "恒发国际", "广州海淘集运", "俄罗斯涅瓦物流",
        "四季安集团", "弘海铁盛", "泛昇国际", "嘉拓智能", "广东晟晖再生资源集团"
    ]
    
    positions = [
        (Inches(1), Inches(2)), (Inches(4.8), Inches(2)), (Inches(8.6), Inches(2)),
        (Inches(1), Inches(3.2)), (Inches(4.8), Inches(3.2)), (Inches(8.6), Inches(3.2)),
        (Inches(1), Inches(4.4)), (Inches(4.8), Inches(4.4)), (Inches(8.6), Inches(4.4)),
    ]
    
    for i, client in enumerate(clients):
        x, y = positions[i]
        
        # 客户卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.5), Inches(0.9))
        card.fill.solid()
        card.fill.fore_color.rgb = PRIMARY_BLUE
        card.line.fill.background()
        
        # 客户名称
        name_box = slide.shapes.add_textbox(x, y + Inches(0.2), Inches(3.5), Inches(0.5))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = client
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_contact_slide(prs):
    """创建联系页/结束页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 渐变蓝色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_BLUE
    bg.line.fill.background()
    
    # 公司名
    company_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.2))
    tf = company_box.text_frame
    p = tf.paragraphs[0]
    p.text = "广州八米网络科技有限公司"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 品牌
    brand_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(12.333), Inches(1))
    tf = brand_box.text_frame
    p = tf.paragraphs[0]
    p.text = "八米科技  BAMIKEJI"
    p.font.size = Pt(36)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # 网址
    web_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.8))
    tf = web_box.text_frame
    p = tf.paragraphs[0]
    p.text = "www.bamitms.com"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 标语
    tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.8))
    tf = tagline_box.text_frame
    p = tf.paragraphs[0]
    p.text = "15年行业经验 · 1200+成功客户 · 让跨境物流更智能"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 底部装饰线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(6.8), Inches(5.333), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()
    
    return slide

# ==================== 开始制作PPT ====================

# 第1页：封面
add_title_slide(prs, 
    "八米科技", 
    "跨境供应链数字化解决方案领导者")

# 第2页：公司简介
add_content_slide(prs, "公司简介", [
    ("公司全称", "广州八米网络科技有限公司"),
    ("品牌名称", "八米科技（BAMIKEJI）"),
    ("官方网站", "www.bamitms.com"),
    ("企业定位", "中国领先的跨境供应链数字化解决方案服务商"),
    ("行业经验", "15年国际物流行业深耕，服务1200+成功客户"),
    ("企业荣誉", "2025年入选广东省科技型中小企业库（登记编号：2025440111A0032458）"),
])

# 第3页：核心业务
add_content_slide(prs, "核心业务定位", [
    "专注为国际物流、集运、海外仓企业提供软件系统支撑",
    "通过数字化手段帮助传统跨境物流企业实现转型升级",
    "构建「获客+集运+代购+商城」一体化管理生态",
    "赋能中国品牌出海，打通端到端供应链业务体系",
])

# 第4页：四大产品
add_product_slide(prs, "四大核心产品", [
    ("代购集运系统", "独家研发的代购商城与集运一体化系统，支持纯集运业务场景。帮助海外华人实现反向海淘，从国内电商平台购买商品转运至海外。"),
    ("跨境TMS系统", "适合海陆空铁运输、集装转运、自营国际专线等各种物流场景。从工厂到客户、跨系统跨部门全链条打通。"),
    ("海外仓WMS系统", "助力海外仓企业搭建企业营销官网、全球化仓库布局，开发代发货、虚拟仓、退换标、中转、调拨等业务。"),
    ("集成供应链", "建设企业面向客户产品交付的端到端业务体系，依据数据驱动、系统支撑，为全球供应链全领域提供数字化协同。"),
])

# 第5页：核心优势
add_advantage_slide(prs)

# 第6页：服务模式
add_content_slide(prs, "服务模式", [
    ("标准化SaaS产品", "开箱即用，按年/月付费，功能成熟，快速上线。适合中小型物流企业，降低数字化门槛。"),
    ("定制化开发服务", "根据企业特殊需求进行深度定制，打造企业专属灵活可拓展的系统。适合大型客户或特殊业务场景。"),
    ("7×24小时客服支持", "拥有强大客服团队，全天候在线为客户答疑解惑，确保系统稳定运行。"),
    ("持续升级迭代", "紧跟行业发展趋势，持续进行产品升级和功能迭代，确保客户始终使用最先进的系统。"),
])

# 第7页：目标客户
add_content_slide(prs, "目标客户群体", [
    "国际货运代理企业",
    "华人集运公司（反向海淘服务商）",
    "自营国际专线物流企业",
    "海外仓运营商（一件代发、FBA中转）",
    "跨境电商平台",
    "品牌出海企业",
    "代购平台及个人代购服务商",
])

# 第8页：客户案例
add_client_slide(prs)

# 第9页：最新动态
add_content_slide(prs, "最新动态", [
    ("2026年3月", "与菲律宾最大零售IT服务商 ASC 签署长期战略合作伙伴协议，联合打造东南亚统仓共配4PL解决方案"),
    ("2025年12月", "成功入选广东省科技型中小企业库，获得科技部门高度认可"),
    ("2025年", "与新能源高端装备领军企业嘉拓智能达成战略合作，进军新能源智能制造物流领域"),
    ("2025年6月", "与广东晟晖再生资源集团合作，建设供应链数智化平台，进军再生资源行业"),
])

# 第10页：结束页
add_contact_slide(prs)

# 保存PPT
output_path = os.path.expanduser("~/Desktop/八米科技企业介绍.pptx")
prs.save(output_path)
print(f"PPT已生成：{output_path}")